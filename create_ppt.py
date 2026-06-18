from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # 슬라이드 레이아웃 설정
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # 1. 표지
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "자율주행 무인차량(UGV)의 사이버 해킹 방어를 위한\n동형암호(FHE) 기반 초경량 이상탐지 시스템"
    subtitle.text = "기말 프로젝트 최종 발표\n발표자: 김동건"

    # 2. 목차
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "발표 순서"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "1. 연구 배경 및 문제 정의"
    p = tf.add_paragraph()
    p.text = "2. 동형암호(FHE) 이론 및 적용"
    p = tf.add_paragraph()
    p.text = "3. 경량화 AI(로지스틱 회귀) 선정 이유"
    p = tf.add_paragraph()
    p.text = "4. 시스템 아키텍처 및 파이프라인"
    p = tf.add_paragraph()
    p.text = "5. 검증 데이터셋 소개 (고려대 HCRL)"
    p = tf.add_paragraph()
    p.text = "6. 최종 모델 검증 결과"
    p = tf.add_paragraph()
    p.text = "7. 결론 및 향후 과제"

    # 3. 연구 배경
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "1. 연구 배경 및 문제 정의"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "현대 군사 작전에서 무인차량(UGV)의 중요성 증대"
    p = tf.add_paragraph()
    p.text = "적군의 악의적인 해킹(속도/조향각 조작, CAN Bus 침입) 위협 존재"
    p = tf.add_paragraph()
    p.text = "기존 방어 시스템의 한계"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "암호화된 데이터를 복호화하여 검사하는 동안 치명적인 지연 시간(Latency) 발생"
    p.level = 2

    # 4. 동형암호 이론
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "2. 동형암호(FHE) 이론 및 적용"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "동형암호(Fully Homomorphic Encryption)란?"
    p = tf.add_paragraph()
    p.text = "데이터를 암호화한 상태 그대로 연산할 수 있는 차세대 암호 기술"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "본 프로젝트의 적용 방안"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "UGV의 주행/센서 데이터를 암호화하여 클라우드로 전송"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "서버는 원본 데이터를 볼 수 없는 상태에서 AI 연산 수행 및 안전 판단"
    p.level = 1

    # 5. 로지스틱 회귀 선정 이유
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "3. 경량화 AI 선정 이유"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "자율주행 방어의 핵심은 '실시간성(Low Latency)' 보장"
    p = tf.add_paragraph()
    p.text = "DNN 등 복잡한 모델의 한계"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "비선형 함수 처리 시 암호문 연산 과부하 발생 (최대 3,400ms)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "초경량 로지스틱 회귀(Logistic Regression) 모델 도입"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "단순 선형 연산으로 오버헤드 최소화 (처리속도 34.4ms 달성)"
    p.level = 2
    
    if os.path.exists('fhe_latency_comparison.png'):
        slide.shapes.add_picture('fhe_latency_comparison.png', Inches(4.5), Inches(2.0), width=Inches(5))

    # 6. 시스템 파이프라인
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "4. 시스템 아키텍처 및 파이프라인"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Client (UGV 차량)"
    p = tf.add_paragraph()
    p.text = "차량 데이터 수집 ➡️ 정규화 ➡️ 공개키 기반 FHE 암호화 ➡️ 서버 전송"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "서버로부터 수신한 암호문 해독 ➡️ 해킹 여부 확인 ➡️ 즉각 제동"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Server (클라우드)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "암호화된 상태로 데이터 수신 ➡️ 로지스틱 회귀 이상탐지 연산 ➡️ 암호문 반환"
    p.level = 1

    # 7. HCRL 데이터셋 소개 (NEW)
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "5. 검증 데이터셋 소개 (고려대 HCRL)"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "고려대학교 HCRL(Hacking and Countermeasure Research Lab)"
    p = tf.add_paragraph()
    p.text = "자율주행 및 자동차 보안 분야에서 국제적으로 가장 널리 인용되는 공인 데이터셋"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "데이터 수집 환경"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "실제 차량(기아 쏘울)의 내부 통신망(CAN Bus)에서 실주행 중 트래픽 추출"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4대 치명적 해킹 시나리오 포함"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "DoS 공격: 통신망 전체 마비 시도"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fuzzy 공격: 무작위 교란 신호 주입"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "RPM/Gear Spoofing 공격: 차량 속도 및 기어 강제 조작 (본 연구의 핵심 시나리오)"
    p.level = 1

    if os.path.exists('hcrl_dataset_architecture.png'):
        slide.shapes.add_picture('hcrl_dataset_architecture.png', Inches(4.5), Inches(2.0), width=Inches(5))


    # 8. 실험 및 검증 결과
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "6. 최종 모델 검증 결과"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "HCRL 데이터셋 대상 로지스틱 회귀 모델의 FHE 검증"
    p = tf.add_paragraph()
    p.text = "탐지 정확도 (Accuracy)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "평문 상태(79.5%)와 동형암호 상태(79.0%)에서 동일한 수준의 높은 방어율 입증"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "초저지연 실시간성 (Ultra-Low Latency)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "암호문 추론 지연시간 평균 2.47ms 달성 (목표치 34.4ms를 월등히 상회)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "결론: 고속 주행 차량에서도 딜레이 없는 완벽한 실시간 해킹 방어 가능"
    p.level = 1

    # 9. 결론 및 향후 과제
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "7. 결론 및 향후 과제"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "결론"
    p = tf.add_paragraph()
    p.text = "동형암호를 통한 데이터 보호와 로지스틱 회귀를 통한 초저지연(2.47ms) 방어 시스템의 융합 성공"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "HCRL 공인 데이터셋을 통해 실전 모델로서의 타당성 및 정확도 입증"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "향후 과제"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "로컬 3D 물리 시뮬레이터(CARLA) 구축 및 엣지 디바이스(젯슨 나노 등) 탑재 테스트 진행"
    p.level = 1

    prs.save('UGV_FHE_Presentation_Final.pptx')
    print("HCRL 소개가 추가된 완성형 PPT 파일 생성 완료!")

if __name__ == '__main__':
    create_presentation()
