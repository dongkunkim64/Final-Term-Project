from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_imrad_presentation():
    prs = Presentation()
    
    # 0. 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "자율주행 UGV 사이버 보안을 위한\nFHE 기반 초저지연 이상탐지 연구"
    slide.placeholders[1].text = "IMRAD 구조 기반 기말 프로젝트 발표\n발표자: 김동건"

    # I. INTRODUCTION
    # 1. 연구 배경
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "[Introduction] 연구 배경 및 문제 정의"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "군사/민간 자율주행 무인차량(UGV)의 보안 위협 증가"
    p = tf.add_paragraph()
    p.text = "차량 내부 통신망(CAN Bus) 해킹 시 속도 및 조향각 원격 조작 위험"
    p = tf.add_paragraph()
    p.text = "기존 클라우드 기반 방어 시스템의 치명적 한계"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "클라우드에서 AI 검사를 수행하려면 데이터를 평문으로 복호화해야 함"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "이 과정에서 심각한 지연시간(Latency) 발생 및 복호화된 원본 데이터 탈취 위험"
    p.level = 2

    # 2. 왜 동형암호인가?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "[Introduction] 왜 '동형암호(FHE)'를 사용했는가?"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "자율주행 시스템 보안의 두 가지 절대 조건: '실시간성'과 '데이터 프라이버시'"
    p = tf.add_paragraph()
    p.text = "동형암호(FHE)의 핵심 가치"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "데이터를 암호화한 상태(Gibberish) 그대로 연산 가능"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "클라우드는 원본 데이터를 전혀 보지 못한 채로, 해킹 여부만 판별하여 암호문으로 반환"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "본 연구 도입의 결정적 이유"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "복호화 단계를 완전히 생략하여 '보안 취약점'을 원천 차단하고 구조를 단순화하기 위함"
    p.level = 2

    # M. METHODS
    # 3. 실험 파이프라인 및 아키텍처
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "[Methods] 제안하는 시스템 파이프라인 (시연)"
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(1.0))
    tf2 = txBox.text_frame
    tf2.text = "1. UGV 데이터 수집 ➔ 2. FHE 암호화 ➔ 3. 클라우드 연산(로지스틱 회귀) ➔ 4. 결과 회신 및 제동"
    
    if os.path.exists('fhe_effect_demonstration.mp4') and os.path.exists('poster.png'):
        slide.shapes.add_movie('fhe_effect_demonstration.mp4', Inches(1.5), Inches(2.2), width=Inches(7), height=Inches(3.94), poster_frame_image='poster.png', mime_type='video/mp4')
    else:
        txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(3.0), Inches(6.0), Inches(1.0))
        tf3 = txBox2.text_frame
        tf3.text = "🎥 영상을 찾을 수 없습니다."

    # 4. 검증 데이터셋 아키텍처
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "[Methods] 실험 데이터셋 아키텍처 (고려대 HCRL)"
    if os.path.exists('hcrl_dataset_architecture.png'):
        slide.shapes.add_picture('hcrl_dataset_architecture.png', Inches(1.5), Inches(1.5), width=Inches(7))

    # R. RESULTS
    # 5. 해킹 공격 탐지 결과
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "[Results] 실시간 RPM Spoofing 공격 탐지 결과"
    if os.path.exists('hcrl_dynamic_attack_graph.gif'):
        slide.shapes.add_picture('hcrl_dynamic_attack_graph.gif', Inches(1.5), Inches(2.0), width=Inches(7))

    # 6. 정확도 및 레이턴시
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "[Results] 모델 최종 성능 검증 (정확도 및 레이턴시)"
    if os.path.exists('fhe_hcrl_validation_animated.gif'):
        slide.shapes.add_picture('fhe_hcrl_validation_animated.gif', Inches(1.0), Inches(2.0), width=Inches(8))
    elif os.path.exists('fhe_hcrl_validation.png'):
        slide.shapes.add_picture('fhe_hcrl_validation.png', Inches(1.0), Inches(2.0), width=Inches(8))

    # D. DISCUSSION
    # 7. 고찰 및 결론
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "[Discussion] 고찰 및 향후 과제"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "연구의 의의 (Discussion)"
    p = tf.add_paragraph()
    p.text = "무겁다고 여겨졌던 동형암호를 초경량 선형 모델(로지스틱 회귀)과 결합 시, 2.47ms라는 초실시간 레이턴시로 자율주행 차량에 적용 가능함을 수학적/실증적으로 증명함."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "평문 검사와 비교해 정확도 손실(0.53%p)이 거의 발생하지 않음을 확인"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "향후 과제 (Future Works)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "본 논리 모델을 기반으로 한 3D 물리 시뮬레이터(CARLA) 정밀 구동 연구"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "소형 엣지 디바이스(NVIDIA Jetson) 기반의 HIL(Hardware-In-the-Loop) 테스트 진행"
    p.level = 1

    prs.save('UGV_FHE_IMRAD_Presentation.pptx')
    print("IMRAD 구조가 적용된 새로운 PPT 'UGV_FHE_IMRAD_Presentation.pptx' 생성 완료!")

if __name__ == '__main__':
    create_imrad_presentation()
