from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_visual_presentation():
    prs = Presentation()
    
    # 1. 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "자율주행 UGV 사이버 해킹 방어\n(FHE 초경량 모델 시연)"
    slide.placeholders[1].text = "고려대학교 HCRL 공인 데이터셋 기반 시각적 검증\n발표자: 김동건"

    # 2. HCRL 데이터셋 아키텍처
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    slide.shapes.title.text = "1. HCRL Car Hacking Dataset 추출 구조"
    if os.path.exists('hcrl_dataset_architecture.png'):
        slide.shapes.add_picture('hcrl_dataset_architecture.png', Inches(1.5), Inches(1.5), width=Inches(7))

    # 3. 해킹 공격 감지 (동적 그래프)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "2. 실시간 RPM Spoofing 해킹 공격 데이터 시각화"
    if os.path.exists('hcrl_dynamic_attack_graph.gif'):
        slide.shapes.add_picture('hcrl_dynamic_attack_graph.gif', Inches(1.5), Inches(2.0), width=Inches(7))

    # 4. FHE 동형암호 방어 메커니즘 시연 영상
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "3. FHE 기반 초저지연 방어 시뮬레이션 (동영상)"
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(3.0), Inches(6.0), Inches(1.0))
    tf = txBox.text_frame
    tf.text = "🎥 여기에 'fhe_effect_demonstration.mp4' 영상을 마우스로 드래그해서 넣어주세요.\n(파워포인트에서 자동 재생 및 전체화면 설정 권장)"

    # 5. FHE 모델 검증 결과 (레이턴시 & 정확도)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "4. 최종 모델 검증: 정확도 및 초저지연(2.47ms)"
    if os.path.exists('fhe_hcrl_validation_animated.gif'):
        slide.shapes.add_picture('fhe_hcrl_validation_animated.gif', Inches(1.0), Inches(2.0), width=Inches(8))

    # 6. 결론
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "결론: HCRL 데이터셋 기반 시연 요약"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "원본 노출 없는 동형암호의 절대적 보안성 시각적 입증"
    p = tf.add_paragraph()
    p.text = "초경량 로지스틱 회귀 모델로 추론 지연시간 2.47ms 달성 (초실시간성 확보)"
    p = tf.add_paragraph()
    p.text = "공인된 차량 해킹(RPM Spoofing) 공격을 즉각적으로 감지하고 방어(Braking) 성공"

    prs.save('UGV_FHE_Visual_Presentation.pptx')
    print("HCRL 시각 자료 위주의 새로운 PPT 'UGV_FHE_Visual_Presentation.pptx' 생성 완료!")

if __name__ == '__main__':
    create_visual_presentation()
